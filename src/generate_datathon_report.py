import os
from pathlib import Path

REPORTS_DIR = Path("reports")
FIGURES_DIR = REPORTS_DIR / "figures"
MPL_CACHE_DIR = REPORTS_DIR / ".matplotlib"

os.environ.setdefault("MPLCONFIGDIR", str(MPL_CACHE_DIR.resolve()))

import matplotlib
import numpy as np
import pandas as pd
import seaborn as sns
from pptx import Presentation
from pptx.util import Inches, Pt

from src.config import OUTPUTS_DIR, TARGET_CANDIDATES, TRAIN_PATH
from src.feature_engineering import add_real_estate_features, add_train_based_location_features
from src.modeling import fit_model, predict_prices, split_features_target
from src.preprocessing import basic_clean, remove_training_outliers
from src.utils import detect_column

matplotlib.use("Agg")

import matplotlib.pyplot as plt


def savefig(name: str) -> None:
    path = FIGURES_DIR / name
    path.parent.mkdir(parents=True, exist_ok=True)
    plt.tight_layout()
    plt.savefig(path, dpi=180, bbox_inches="tight")
    plt.close()


def add_price_per_sqft(df: pd.DataFrame, target_col: str) -> pd.DataFrame:
    df = df.copy()
    if "total_sqft_clean" in df.columns:
        df["price_per_sqft"] = df[target_col] / df["total_sqft_clean"].replace(0, np.nan)
    elif "area" in df.columns:
        df["price_per_sqft"] = df[target_col] / df["area"].replace(0, np.nan)
    return df


def plot_missing_values(df: pd.DataFrame) -> None:
    missing_pct = df.isna().mean().sort_values(ascending=False) * 100
    plt.figure(figsize=(10, 5))
    sns.barplot(x=missing_pct.index, y=missing_pct.values, color="#4C78A8")
    plt.xticks(rotation=45, ha="right")
    plt.ylabel("Missing values (%)")
    plt.title("Missing Values By Column")
    savefig("01_missing_values.png")


def plot_price_distribution(df: pd.DataFrame, target_col: str) -> None:
    plt.figure(figsize=(10, 5))
    sns.histplot(df[target_col], kde=True, bins=35, color="#F58518")
    plt.xlabel("Price (INR)")
    plt.title("Target Price Distribution")
    savefig("02_price_distribution.png")

    plt.figure(figsize=(10, 5))
    sns.histplot(np.log1p(df[target_col]), kde=True, bins=35, color="#54A24B")
    plt.xlabel("Log transformed price")
    plt.title("Log Price Distribution")
    savefig("03_log_price_distribution.png")


def plot_numeric_correlation(df: pd.DataFrame) -> None:
    numeric = df.select_dtypes(include=[np.number])
    corr = numeric.corr()
    plt.figure(figsize=(11, 8))
    sns.heatmap(corr, cmap="vlag", center=0, linewidths=0.4)
    plt.title("Numeric Feature Correlation Heatmap")
    savefig("04_numeric_correlation_heatmap.png")


def plot_price_vs_area(df: pd.DataFrame, target_col: str) -> None:
    area_col = "total_sqft_clean" if "total_sqft_clean" in df.columns else "area"
    plt.figure(figsize=(10, 6))
    sns.scatterplot(data=df, x=area_col, y=target_col, hue="bhk" if "bhk" in df.columns else None, alpha=0.65)
    plt.xlabel("Area (sqft)")
    plt.ylabel("Price (INR)")
    plt.title("Price vs Area")
    savefig("05_price_vs_area.png")


def plot_location_signals(df: pd.DataFrame, target_col: str) -> None:
    if "location_clean" not in df.columns:
        return

    location_price = (
        df.groupby("location_clean")[target_col]
        .agg(["count", "median"])
        .query("count >= 5")
        .sort_values("median", ascending=False)
        .head(20)
        .reset_index()
    )
    plt.figure(figsize=(11, 7))
    sns.barplot(data=location_price, y="location_clean", x="median", color="#E45756")
    plt.xlabel("Median price (INR)")
    plt.ylabel("Location")
    plt.title("Top Locations By Median Price")
    savefig("06_top_locations_by_median_price.png")

    if "price_per_sqft" in df.columns:
        location_pps = (
            df.groupby("location_clean")["price_per_sqft"]
            .median()
            .sort_values(ascending=False)
            .head(20)
            .reset_index()
        )
        plt.figure(figsize=(11, 7))
        sns.barplot(data=location_pps, y="location_clean", x="price_per_sqft", color="#72B7B2")
        plt.xlabel("Median price per sqft")
        plt.ylabel("Location")
        plt.title("Location Premium Signal: Median Price Per Sqft")
        savefig("07_location_price_per_sqft.png")


def plot_category_effects(df: pd.DataFrame, target_col: str) -> None:
    columns = [col for col in ["furnishing", "property_type", "bhk", "bath"] if col in df.columns]
    if not columns:
        return

    fig, axes = plt.subplots(2, 2, figsize=(13, 9))
    axes = axes.flatten()
    for ax, col in zip(axes, columns):
        order = df.groupby(col)[target_col].median().sort_values(ascending=False).index
        sns.boxplot(data=df, x=col, y=target_col, order=order, ax=ax)
        ax.set_title(f"Price By {col}")
        ax.tick_params(axis="x", rotation=35)
    for ax in axes[len(columns) :]:
        ax.axis("off")
    savefig("08_category_price_patterns.png")


def plot_model_comparison() -> None:
    comparison_path = OUTPUTS_DIR / "model_comparison.csv"
    if not comparison_path.exists():
        return
    scores = pd.read_csv(comparison_path)
    plt.figure(figsize=(9, 5))
    sns.barplot(data=scores, x="model", y="mean_rmsle", color="#B279A2")
    plt.ylabel("Mean RMSLE")
    plt.xlabel("Model")
    plt.title("Model Comparison: Lower RMSLE Is Better")
    savefig("09_model_comparison.png")


def plot_model_diagnostics(df: pd.DataFrame, target_col: str, best_model_name: str) -> None:
    df = df.drop(columns=["price_per_sqft"], errors="ignore")
    X, y = split_features_target(df, target_col)
    model = fit_model(X, y, model_name=best_model_name)

    preds = predict_prices(model, X)
    residuals = y - preds

    plt.figure(figsize=(7, 7))
    sns.scatterplot(x=y, y=preds, alpha=0.65)
    low = min(y.min(), preds.min())
    high = max(y.max(), preds.max())
    plt.plot([low, high], [low, high], color="#E45756", linestyle="--")
    plt.xlabel("Actual price (INR)")
    plt.ylabel("Predicted price (INR)")
    plt.title("Actual vs Predicted Prices")
    savefig("10_actual_vs_predicted.png")

    plt.figure(figsize=(10, 5))
    sns.scatterplot(x=preds, y=residuals, alpha=0.65)
    plt.axhline(0, color="#E45756", linestyle="--")
    plt.xlabel("Predicted price (INR)")
    plt.ylabel("Residual: actual - predicted")
    plt.title("Residual Pattern Check")
    savefig("11_residuals.png")

    estimator = model.named_steps["model"]
    if hasattr(estimator, "coef_"):
        names = model.named_steps["preprocess"].get_feature_names_out()
        importance = (
            pd.DataFrame({"feature": names, "importance": np.abs(estimator.coef_)})
            .sort_values("importance", ascending=False)
            .head(20)
        )
        plt.figure(figsize=(11, 7))
        sns.barplot(data=importance, y="feature", x="importance", color="#4C78A8")
        plt.xlabel("Absolute coefficient on log price")
        plt.ylabel("Feature")
        plt.title("Top Ridge Model Signals")
        savefig("12_feature_importance.png")


def write_reports(raw: pd.DataFrame, cleaned: pd.DataFrame, target_col: str, best_model_name: str) -> None:
    REPORTS_DIR.mkdir(parents=True, exist_ok=True)
    raw_rows = len(raw)
    clean_rows = len(cleaned)
    removed_rows = raw_rows - clean_rows
    location_count = raw["location"].nunique() if "location" in raw.columns else "N/A"

    eda = f"""# DataVerse 2 PM Technical EDA Report

## Problem

Predict Bengaluru residential property prices from structured property attributes. This is a regression problem, and the contest target metric for the housing track is RMSLE.

## Dataset Snapshot

- Raw rows: {raw_rows}
- Rows after cleaning/outlier removal: {clean_rows}
- Removed outlier/noisy rows: {removed_rows}
- Columns: {len(raw.columns)}
- Unique locations: {location_count}
- Target column: `{target_col}`

## Cleaning Decisions

- Standardized location strings into a lowercase `location_clean` field.
- Converted area into a numeric sqft feature.
- Removed invalid or unrealistic training records, including negative prices, impossible BHK/bath counts, and unusual sqft-per-BHK values.
- Preserved categorical columns for one-hot encoding inside the model pipeline.

## Hidden Signals Found

- Location is a strong price signal, especially through median price per sqft.
- Area and BHK are useful, but area alone is not enough because premium localities can dominate price.
- Engineered ratios such as sqft per BHK and bath per BHK help describe compact versus spacious properties.
- Furnishing, parking, balcony, and property type add smaller but useful signals.

## Generated Plots

![Missing values](figures/01_missing_values.png)
![Price distribution](figures/02_price_distribution.png)
![Log price distribution](figures/03_log_price_distribution.png)
![Correlation heatmap](figures/04_numeric_correlation_heatmap.png)
![Price vs area](figures/05_price_vs_area.png)
![Top locations](figures/06_top_locations_by_median_price.png)
![Location price per sqft](figures/07_location_price_per_sqft.png)
![Category price patterns](figures/08_category_price_patterns.png)
![Model comparison](figures/09_model_comparison.png)
![Actual vs predicted](figures/10_actual_vs_predicted.png)
![Residuals](figures/11_residuals.png)
![Feature importance](figures/12_feature_importance.png)
"""
    (REPORTS_DIR / "DATATHON_EDA_REPORT.md").write_text(eda, encoding="utf-8")

    strategy = f"""# DataVerse Winning Strategy

## What The Contest Requires

The guide and presentation describe a 10-hour datathon with four tracks. For the Bengaluru housing track, the goal is to estimate residential property prices and minimize RMSLE. Judges care about:

- Model rigor: algorithm choice and hyperparameter justification.
- Technical pivot: ability to merge new helper data released at 3 PM.
- Stability: consistent performance on hidden/private test data.
- Interpretability: strong EDA, feature importance, and clear plots.
- Pitch quality: explaining why predictions make sense.

## What This Project Already Covers

- Reproducible training pipeline in `src/run_pipeline.py`.
- Cleaning and outlier handling in `src/preprocessing.py`.
- Feature engineering in `src/feature_engineering.py`.
- Model comparison in `src/modeling.py`.
- Best current model: `{best_model_name}`.
- Streamlit demo app in `app.py`.
- EDA and model plots in `reports/figures/`.
- 2 PM audit report in `reports/DATATHON_EDA_REPORT.md`.

## What To Do To Maximize Winning Chances

1. Prepare the 2 PM EDA pitch around cleaning, hidden signals, and plots.
2. Be ready for the 3 PM pivot by writing new helper data as CSV into `data/` and merging it by location or property identifiers.
3. Track before/after metric changes after adding pivot data.
4. Add ablation evidence: compare baseline features versus engineered features.
5. Keep the final GitHub repo clean and freeze it before the deadline.

## Strong Judge Explanation

We used RMSLE because housing prices vary widely, and log error is more stable for price prediction. We tested Ridge, Random Forest, Extra Trees, and XGBoost, then selected the best performer by cross-validation. The main insight is that location premium and engineered space-quality ratios explain more than raw area alone.
"""
    (REPORTS_DIR / "DATATHON_WINNING_STRATEGY.md").write_text(strategy, encoding="utf-8")

    script = """# DataVerse Presentation Script

## 1. Problem

Our track is Bengaluru house price prediction. The task is to predict residential property prices from property features while minimizing RMSLE.

## 2. Data Understanding

The dataset contains property attributes such as area, location, BHK, bathrooms, balcony, parking, furnishing, property type, age, and price.

## 3. Cleaning

We standardized locations, converted area to numeric sqft, checked missing values, and removed unrealistic outliers such as impossible BHK/bathroom counts or abnormal sqft-per-BHK values.

## 4. EDA Insights

Price is skewed, so log transformation helps. Location is a major signal. Area matters, but premium locations can cause large price differences even for similar-sized homes.

## 5. Feature Engineering

We added sqft per BHK, bath per BHK, balcony and parking flags, furnishing flags, property type flags, premium location keyword flags, and location-level median price-per-sqft features.

## 6. Modeling

We compared Ridge Regression, Random Forest, Extra Trees, and XGBoost with 5-fold cross-validation. Ridge Regression currently gives the best RMSLE, so we use it as the saved final model.

## 7. Limitations

The model does not know exact coordinates, metro distance, road access, school quality, or market trends. These would be useful pivot/helper features if provided.

## 8. Impact

The app gives a quick market estimate and confidence range, useful for screening properties and negotiation support.
"""
    (REPORTS_DIR / "PRESENTATION_SCRIPT.md").write_text(script, encoding="utf-8")


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.title
    title_box.text = title
    title_box.text_frame.paragraphs[0].font.size = Pt(32)
    if subtitle:
        box = slide.shapes.add_textbox(Inches(0.7), Inches(1.7), Inches(12), Inches(1.1))
        frame = box.text_frame
        frame.text = subtitle
        frame.paragraphs[0].font.size = Pt(18)


def add_bullets(slide, bullets: list[str], top: float = 1.7) -> None:
    box = slide.shapes.add_textbox(Inches(0.8), Inches(top), Inches(12), Inches(4.8))
    frame = box.text_frame
    frame.clear()
    for idx, bullet in enumerate(bullets):
        paragraph = frame.paragraphs[0] if idx == 0 else frame.add_paragraph()
        paragraph.text = bullet
        paragraph.level = 0
        paragraph.font.size = Pt(18)


def add_image_slide(prs: Presentation, title: str, image_name: str, note: str) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, title)
    image_path = FIGURES_DIR / image_name
    slide.shapes.add_picture(str(image_path), Inches(0.9), Inches(1.35), width=Inches(7.5))
    box = slide.shapes.add_textbox(Inches(8.7), Inches(1.7), Inches(4.1), Inches(4.6))
    frame = box.text_frame
    frame.text = note
    frame.paragraphs[0].font.size = Pt(17)


def build_pitch_deck(raw: pd.DataFrame, cleaned: pd.DataFrame, best_model_name: str) -> None:
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    slide = prs.slides.add_slide(prs.slide_layouts[0])
    slide.shapes.title.text = "Bengaluru House Price Prediction"
    slide.placeholders[1].text = "DataVerse Housing Track | Regression | RMSLE"

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "Contest Requirement")
    add_bullets(
        slide,
        [
            "Build a regression model for Bengaluru residential property prices.",
            "Optimize for RMSLE because the target price range is wide and skewed.",
            "Show strong EDA, cleaning decisions, feature engineering, model rigor, and interpretability.",
            "Be ready for the 3 PM technical pivot by merging new helper data and proving metric improvement.",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "Dataset And Cleaning")
    add_bullets(
        slide,
        [
            f"Raw rows: {len(raw)}; cleaned rows: {len(cleaned)}; removed outlier/noisy rows: {len(raw) - len(cleaned)}.",
            "Cleaned location strings and converted area into numeric sqft.",
            "Removed impossible or unrealistic BHK, bathroom, price, and sqft-per-BHK records.",
            "Preserved categorical features for one-hot encoding in the model pipeline.",
        ],
    )

    add_image_slide(
        prs,
        "Target Distribution",
        "03_log_price_distribution.png",
        "House prices are skewed. Log transformation makes model training and RMSLE evaluation more stable.",
    )
    add_image_slide(
        prs,
        "Location Is The Strongest Signal",
        "07_location_price_per_sqft.png",
        "Median price per sqft by location captures locality premium better than raw area alone.",
    )
    add_image_slide(
        prs,
        "Area Matters, But Not Alone",
        "05_price_vs_area.png",
        "Area and BHK explain part of price, but similar-sized homes can vary strongly by location and property profile.",
    )
    add_image_slide(
        prs,
        "Model Comparison",
        "09_model_comparison.png",
        f"The best cross-validated model is {best_model_name}. Lower RMSLE means better price prediction stability.",
    )
    add_image_slide(
        prs,
        "Prediction Diagnostics",
        "10_actual_vs_predicted.png",
        "Actual vs predicted values show whether the model follows the correct pricing trend across the data.",
    )
    add_image_slide(
        prs,
        "Interpretability",
        "12_feature_importance.png",
        "Ridge coefficients show which engineered and encoded features most influence log-price prediction.",
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "3 PM Pivot Readiness")
    add_bullets(
        slide,
        [
            "If helper data is released, merge it by location or property identifier.",
            "Possible helper features: metro distance, IT park proximity, school ratings, road access, zoning, or guidance value.",
            "Create a before/after metric table to prove the delta after adding pivot features.",
            "Explain whether the provided hypothesis was supported or disproved by plots and validation score.",
        ],
    )

    slide = prs.slides.add_slide(prs.slide_layouts[5])
    add_title(slide, "Final Pitch")
    add_bullets(
        slide,
        [
            "This project covers the full ML workflow: EDA, cleaning, feature engineering, validation, saved model, and app.",
            "Main insight: location premium plus space-quality ratios explain price better than area alone.",
            "Main limitation: exact coordinates, infrastructure proximity, and market trend data are not present.",
            "Impact: quick market screening and negotiation baseline for Bengaluru properties.",
        ],
    )

    prs.save(REPORTS_DIR / "DataVerse_Housing_Pitch.pptx")


def main() -> None:
    sns.set_theme(style="whitegrid")
    FIGURES_DIR.mkdir(parents=True, exist_ok=True)

    raw = pd.read_csv(TRAIN_PATH)
    target_col = detect_column(raw.columns, TARGET_CANDIDATES)
    if target_col is None:
        raise ValueError(f"Could not detect target column from {TARGET_CANDIDATES}.")

    cleaned = basic_clean(raw)
    cleaned = remove_training_outliers(cleaned, target_col)
    cleaned = add_real_estate_features(cleaned)
    cleaned, _ = add_train_based_location_features(cleaned, cleaned.copy(), target_col)
    cleaned = add_price_per_sqft(cleaned, target_col)

    scores_path = OUTPUTS_DIR / "model_comparison.csv"
    best_model_name = "ridge"
    if scores_path.exists():
        best_model_name = pd.read_csv(scores_path).iloc[0]["model"]

    plot_missing_values(raw)
    plot_price_distribution(cleaned, target_col)
    plot_numeric_correlation(cleaned)
    plot_price_vs_area(cleaned, target_col)
    plot_location_signals(cleaned, target_col)
    plot_category_effects(cleaned, target_col)
    plot_model_comparison()
    plot_model_diagnostics(cleaned, target_col, best_model_name)
    write_reports(raw, cleaned, target_col, best_model_name)
    build_pitch_deck(raw, cleaned, best_model_name)

    print(f"Generated datathon reports in {REPORTS_DIR.resolve()}")


if __name__ == "__main__":
    main()
